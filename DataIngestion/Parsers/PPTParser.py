from pptx import Presentation
from exception import DocumentParsingError

class PPTParser:
    def __init__(self, ppt_path):
        self.ppt_path = ppt_path
        try:
            self.ppt = Presentation(self.ppt_path)
        
        except Exception as e:
            raise DocumentParsingError("Error loading ppt file")
        
        
    # function to parse the ppt
    def parse(self):
        content = ""
        try:
            for slide_num, slide in enumerate(self.ppt.slides, start=1):
                content += f"[SLIDE {slide_num}] "
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text
                            cleaned_text = " ".join(text.split())
                            content += cleaned_text + " "
        except Exception as e:
            raise DocumentParsingError("Error parsing PPT")

        return {
            "file_name": self.ppt_path,
            "file_type": "pptx",
            "content": content
        }